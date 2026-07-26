from __future__ import annotations

import csv
import json as json_lib
import logging

import fitz  # PyMuPDF
from docx import Document


def extract_text(file_path: str, file_type: str) -> str:
    try:
        if file_type == "pdf":
            doc = fitz.open(file_path)
            text = ""
            try:
                for page in doc:
                    text += page.get_text()
            except TypeError:
                try:
                    for i in range(doc.page_count):
                        page = doc.load_page(i)
                        text += page.get_text()
                except Exception:
                    return ""
            return text
        elif file_type == "docx":
            doc = Document(file_path)
            parts = []
            for para in doc.paragraphs:
                parts.append(para.text)
            return "\n".join(parts) + ("\n" if parts else "")
        elif file_type == "xlsx" or file_type == "xls":
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, data_only=True)
                text = ""
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    text += f"Sheet: {sheet}\n"
                    for row in ws.iter_rows(values_only=True):
                        text += " | ".join(str(cell) if cell is not None else "" for cell in row) + "\n"
                    text += "\n"
                return text
            except ImportError:
                return ""
        elif file_type == "image":
            return ""
        elif file_type == "txt":
            with open(file_path, encoding='utf-8') as f:
                return f.read()
        elif file_type == "pptx":
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                parts = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text)
                    if slide_texts:
                        parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))
                return "\n\n".join(parts)
            except ImportError:
                return ""
        elif file_type == "csv":
            rows = []
            with open(file_path, encoding='utf-8') as f:
                sample = f.read(2048)
                f.seek(0)
                dialect = csv.Sniffer().sniff(sample) if sample else csv.excel
                reader = csv.DictReader(f, dialect=dialect)
                if reader.fieldnames:
                    rows.append(" | ".join(reader.fieldnames))
                    rows.append("-" * len(" | ".join(reader.fieldnames)))
                    for row in reader:
                        rows.append(" | ".join(row.get(col, "") for col in reader.fieldnames))
            return "\n".join(rows) if rows else ""
        elif file_type == "json":
            with open(file_path, encoding='utf-8') as f:
                data = json_lib.load(f)
            return json_lib.dumps(data, indent=2)
        else:
            return ""
    except Exception:
        return ""


def extract_text_pages(file_path: str, file_type: str) -> list[str]:
    """Extract text per page for PDFs, returns list of page texts."""
    if file_type == "pdf":
        try:
            doc = fitz.open(file_path)
            pages = []
            try:
                for page in doc:
                    pages.append(page.get_text())
            except TypeError:
                for i in range(doc.page_count):
                    page = doc.load_page(i)
                    pages.append(page.get_text())
            doc.close()
            return pages
        except Exception:
            return []
    return []


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        if end >= len(text):
            chunks.append(text[start:])
            break
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return chunks


def extract_text_ocr(filepath: str, file_type: str) -> str:
    try:
        import pytesseract
    except ImportError:
        logging.warning("OCR not available. Install pytesseract and Tesseract binary for OCR support.")
        return ""

    try:
        if file_type in ("png", "jpg", "jpeg", "gif", "bmp", "tiff", "image"):
            from PIL import Image
            return pytesseract.image_to_string(Image.open(filepath))
        elif file_type == "pdf":
            doc = fitz.open(filepath)
            text_parts = []
            for page_num in range(doc.page_count):
                page = doc.load_page(page_num)
                pix = page.get_pixmap()
                img_data = pix.tobytes("png")
                import io

                from PIL import Image
                img = Image.open(io.BytesIO(img_data))
                text_parts.append(pytesseract.image_to_string(img))
            doc.close()
            return "\n".join(text_parts)
        return ""
    except Exception:
        return ""
